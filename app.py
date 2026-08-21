import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
import streamlit as st

st.set_page_config(
    page_title="Dynamic Font & Color PPT Automator",
    page_icon="📊",
    layout="centered",
)

st.title("📊 THIS TOOL TO CONVERT DIMENSIONS SIZE FROM INCHES TO FEET OR FEET TO INCHES IN PPT")
st.write(
    "Yeh script slide me maujood neighbouring boxes se **Font Size, Color aur Font Family automatically detect** karke same style apply karegi."
)

st.markdown("---")

uploaded_excel = st.file_uploader(
    "1. Upload Excel File (.xlsx)", type=["xlsx"]
)
uploaded_ppt = st.file_uploader("2. Upload PPT File (.pptx)", type=["pptx"])

if uploaded_excel and uploaded_ppt:
    if st.button("🚀 Process & Regenerate PPT Size Box", type="primary"):
        with st.spinner("Detecting Dynamic Font Sizes & Processing..."):
            try:
                # 1. Excel Read safely
                xl_file = pd.ExcelFile(uploaded_excel)
                sheet_to_use = (
                    "Merged_Result"
                    if "Merged_Result" in xl_file.sheet_names
                    else xl_file.sheet_names[0]
                )
                df = pd.read_excel(uploaded_excel, sheet_name=sheet_to_use)

                # 2. Find Size Column
                size_column_name = None
                for col in df.columns:
                    if "size" in str(col).lower():
                        size_column_name = col
                        break

                prs = Presentation(uploaded_ppt)

                # Process Slides
                for i, slide in enumerate(prs.slides):
                    try:
                        if size_column_name:
                            size_val = str(df[size_column_name].iloc[i]).strip()
                        else:
                            size_val = str(df.iloc[i, 7]).strip()
                    except Exception:
                        size_val = ""

                    # Dynamic Defaults (Fallback agar kisi shape me properties na mile)
                    border_color = RGBColor(227, 108, 10)
                    line_width = Pt(2.0)
                    font_name = "Calibri"
                    font_color = RGBColor(0, 0, 0)
                    detected_font_size = Pt(16)  # Default fallback size

                    left_anchor = None
                    right_anchor = None

                    # STEP 1: DETECT ANCHOR BOXES DYNAMICALLY
                    bottom_shapes = []
                    for shape in slide.shapes:
                        if shape.top > Pt(350) and shape.has_text_frame:
                            txt = shape.text_frame.text.strip().lower()
                            if any(
                                k in txt
                                for k in [
                                    "media",
                                    "type",
                                    "qty",
                                    "remark",
                                    "outlet",
                                    "address",
                                ]
                            ):
                                bottom_shapes.append(shape)

                    bottom_shapes.sort(key=lambda s: s.left)

                    if len(bottom_shapes) >= 2:
                        left_anchor = bottom_shapes[0]
                        right_anchor = bottom_shapes[1]
                    elif len(bottom_shapes) == 1:
                        left_anchor = bottom_shapes[0]

                    # STEP 2: AUTO-DETECT & CLONE FONT SIZE, COLOR, NAME & BORDER
                    ref_shape = left_anchor or right_anchor
                    if ref_shape:
                        # Extract Border Line Style
                        try:
                            if ref_shape.line and ref_shape.line.fill.type == 1:
                                border_color = ref_shape.line.color.rgb
                            if ref_shape.line and ref_shape.line.width:
                                line_width = ref_shape.line.width
                        except Exception:
                            pass

                        # Deep Inspect Paragraphs & Runs for Exact Font Size & Color
                        try:
                            if (
                                ref_shape.has_text_frame
                                and ref_shape.text_frame.paragraphs
                            ):
                                for p in ref_shape.text_frame.paragraphs:
                                    for run in p.runs:
                                        # Auto-Detect Font Size
                                        if run.font and run.font.size:
                                            detected_font_size = run.font.size
                                        # Auto-Detect Font Name
                                        if run.font and run.font.name:
                                            font_name = run.font.name
                                        # Auto-Detect Font Color
                                        if (
                                            run.font
                                            and run.font.color
                                            and run.font.color.rgb
                                        ):
                                            font_color = run.font.color.rgb
                        except Exception:
                            pass

                    # STEP 3: CALCULATE BOUNDARIES DYNAMICALLY
                    if left_anchor and right_anchor:
                        left_bound = left_anchor.left + left_anchor.width
                        right_bound = right_anchor.left
                        top_bound = min(left_anchor.top, right_anchor.top) - Pt(20)
                        bottom_bound = (
                            max(
                                left_anchor.top + left_anchor.height,
                                right_anchor.top + right_anchor.height,
                            )
                            + Pt(20)
                        )
                        target_top = left_anchor.top
                        target_height = left_anchor.height
                    elif left_anchor:
                        left_bound = left_anchor.left + left_anchor.width
                        right_bound = left_anchor.left + left_anchor.width + Pt(140)
                        top_bound = left_anchor.top - Pt(20)
                        bottom_bound = left_anchor.top + left_anchor.height + Pt(20)
                        target_top = left_anchor.top
                        target_height = left_anchor.height
                    else:
                        left_bound = Pt(220)
                        right_bound = Pt(410)
                        top_bound = Pt(380)
                        bottom_bound = Pt(500)
                        target_top = Pt(432)
                        target_height = Pt(28)

                    shapes_to_delete = []

                    # STEP 4: SWEEP & PURGE GHOST TEXTS IN MIDDLE GAP
                    for shape in slide.shapes:
                        if shape == left_anchor or shape == right_anchor:
                            continue

                        txt = ""
                        if shape.has_text_frame:
                            txt = shape.text_frame.text.strip()
                        txt_lower = txt.lower()

                        if any(
                            k in txt_lower
                            for k in [
                                "close view",
                                "far view",
                                "media type",
                                "type:",
                                "qty:",
                                "remarks:",
                            ]
                        ):
                            continue

                        is_in_gap_x = (shape.left >= (left_bound - Pt(15))) and (
                            (shape.left + shape.width) <= (right_bound + Pt(15))
                        )
                        is_in_gap_y = (shape.top >= top_bound) and (
                            (shape.top + shape.height) <= bottom_bound
                        )

                        is_size_label = "size" in txt_lower or (
                            "x" in txt_lower and len(txt) < 30
                        )

                        if (is_in_gap_x and is_in_gap_y) or is_size_label:
                            shapes_to_delete.append(shape)

                    # HARD XML DELETION
                    for shape in shapes_to_delete:
                        try:
                            if shape.has_text_frame:
                                shape.text_frame.text = ""
                            sp = shape._element
                            sp.getparent().remove(sp)
                        except Exception:
                            pass

                    # STEP 5: CREATE RECTANGLE SIZE BOX WITH DETECTED STYLES
                    if size_val and size_val.lower() != "nan":
                        final_text = (
                            size_val
                            if size_val.lower().startswith("size")
                            else f"Size: {size_val}"
                        )

                        box_left = left_bound + Pt(8)
                        box_width = (right_bound - left_bound) - Pt(16)

                        if box_width < Pt(70):
                            box_width = Pt(130)

                        new_box = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            box_left,
                            target_top,
                            box_width,
                            target_height,
                        )

                        new_box.fill.background()
                        new_box.line.color.rgb = border_color
                        new_box.line.width = line_width

                        tf = new_box.text_frame
                        tf.word_wrap = False
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        tf.margin_top = Pt(1)
                        tf.margin_bottom = Pt(1)
                        tf.margin_left = Pt(1)
                        tf.margin_right = Pt(1)

                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER

                        run = p.add_run()
                        run.text = final_text
                        run.font.name = font_name
                        run.font.bold = True  # Always Bold

                        # Apply Auto-Detected Font Color
                        try:
                            run.font.color.rgb = font_color
                        except Exception:
                            pass

                        # Apply Auto-Detected Font Size
                        run.font.size = detected_font_size

                output_ppt_path = "Updated_Presentation.pptx"
                prs.save(output_ppt_path)

                st.success(
                    "🎉 Auto-Detection Complete! Slide se Exact Font Size, Color, aur Font Family detect karke Size Box par apply ho chuka hai."
                )

                with open(output_ppt_path, "rb") as file:
                    st.download_button(
                        label="📥 Download Updated PPT",
                        data=file,
                        file_name="Updated_Presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

            except Exception as e:
                st.error(f"❌ Error Aaya: {e}")
