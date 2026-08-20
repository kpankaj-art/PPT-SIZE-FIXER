import os
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# Relative Paths (Repo ke files)
excel_path = "input_excel.xlsx"
ppt_path = "input_ppt.pptx"
output_ppt_path = "Updated_Presentation.pptx"

if not os.path.exists(excel_path) or not os.path.exists(ppt_path):
    print("❌ Error: Input Excel ya PPT file repo me nahi mili!")
    exit(1)

# Excel sheet read karein
df = pd.read_excel(excel_path, sheet_name="Merged_Result")

# ==========================================
# AUTOMATIC COLUMN SEARCH LOGIC (Excel)
# ==========================================
size_column_name = None

# Excel ke sabhi column names me 'size' keyword dhoondho
for col in df.columns:
    if "size" in str(col).lower():
        size_column_name = col
        break

if size_column_name:
    print(f"✅ Excel me Automatic Size Column mil gaya: '{size_column_name}'")
else:
    print(
        "⚠️ Warning: Excel me 'Size' naam ka column nahi mila! Fallback Column Index 7 (Column H) use kar rahe hain."
    )

prs = Presentation(ppt_path)

last_left = None
last_top = None
last_width = None
last_height = None
last_color = None

manual_check_slides = []

for i, slide in enumerate(prs.slides):
    slide_no = i + 1

    # Dynamic Column Value Read Logic
    try:
        if size_column_name:
            size_val = str(df[size_column_name].iloc[i]).strip()
        else:
            size_val = str(df.iloc[i, 7]).strip()  # Fallback to Column H
    except Exception:
        size_val = ""

    box_found = False
    shapes_to_remove = []

    # STEP 1: Current Slide par Purana Box Search Karein
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "Size" in shape.text_frame.text:
                shapes_to_remove.append(shape)
                box_found = True

                last_left = shape.left
                last_top = shape.top
                last_width = shape.width
                last_height = shape.height

                try:
                    if shape.line.fill.type == 1:
                        last_color = shape.line.color.rgb
                except Exception:
                    pass

    # STEP 2: Delete Old Box
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # STEP 3: Create New Box (If Reference Exists)
    if size_val and size_val.lower() != "nan":
        if last_left is not None:
            new_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                last_left,
                last_top,
                last_width,
                last_height,
            )

            new_box.fill.background()

            if last_color is None:
                last_color = RGBColor(227, 108, 10)

            new_box.line.color.rgb = last_color
            new_box.line.width = Pt(3.0)

            tf = new_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"Size: {size_val}"
            p.alignment = PP_ALIGN.CENTER

            run = p.runs[0]
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 80, 160)
        else:
            manual_check_slides.append(slide_no)

# Save Processed PPT
prs.save(output_ppt_path)

print("\n" + "=" * 50)
print("PROCESS COMPLETED SUCCESSFULLY!")
print("=" * 50)

if manual_check_slides:
    print(
        f"\n⚠️ WARNING: Niche di gayi Slides par koi box/reference nahi mila:"
    )
    print(f"Slide Numbers: {manual_check_slides}")
    print("👉 In slides par Size Box ko MANUALLY check karein.")
else:
    print("\n✅ Sabhi slides successfully update ho gayi hain!")
print("=" * 50)
